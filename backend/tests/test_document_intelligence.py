from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

import pytest

from app.services import document_intelligence_service as svc
from app.tools import document_tools
from app.tools.registry import ToolRegistry


@pytest.fixture(autouse=True)
def _force_builtin_parser(monkeypatch):
    """Make these contract tests deterministic across environments.

    CI installs docling/unstructured from the resolved lock while developer
    machines may not have them; both must exercise the builtin fallback
    semantics asserted below.
    """

    def _unavailable(_path: Path):
        raise svc.AdvancedParserUnavailable("heavy parser disabled for builtin-contract tests")

    monkeypatch.setattr(svc, "_parse_with_docling", _unavailable)
    monkeypatch.setattr(svc, "_parse_with_unstructured", _unavailable)


class _StubProvider:
    name = "stub"

    def __init__(self, reply: str | None = None) -> None:
        self.reply = reply or ""
        self.calls: list[list[dict[str, Any]]] = []

    async def chat(self, messages, model=None, temperature=None, tools=None) -> str:  # noqa: ANN001, ARG002
        self.calls.append(messages)
        return self.reply


def test_parse_advanced_txt_uses_builtin_fallback_without_heavy_deps(tmp_path: Path):
    path = tmp_path / "memo.txt"
    path.write_text("Executive summary\n\nRevenue grew 12% in Q1.", encoding="utf-8")

    ir = svc.parse_advanced(path)

    assert ir.path == str(path)
    assert ir.kind == "text"
    assert ir.parse_engine == "builtin"
    assert ir.pages
    assert len(ir.blocks) == 2
    assert "Revenue grew" in ir.text
    assert ir.document_id.startswith(("blake3:", "sha256:"))


def test_parse_advanced_logs_unexpected_parser_failures(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    caplog: pytest.LogCaptureFixture,
):
    path = tmp_path / "memo.txt"
    path.write_text("Executive summary", encoding="utf-8")

    def _failed_parser(_path: Path):
        raise RuntimeError("parser failed token=supersecrettokenvalue1234567890")

    monkeypatch.setattr(svc, "_parse_with_docling", _failed_parser)
    caplog.set_level(logging.WARNING, logger="app.services.document_intelligence_service")

    ir = svc.parse_advanced(path)

    assert ir.parse_engine == "builtin"
    assert any("docling parser failed" in warning for warning in ir.warnings)
    assert "supersecrettokenvalue" not in "\n".join(ir.warnings)
    assert "document_intelligence.advanced_parser" in caplog.text
    assert "supersecrettokenvalue" not in caplog.text


def test_extract_tables_from_csv(tmp_path: Path):
    path = tmp_path / "sales.csv"
    path.write_text("Region,Revenue\nNorth,100\nSouth,120\n", encoding="utf-8")

    result = svc.extract_tables(path)

    assert result["tables"][0]["headers"] == ["Region", "Revenue"]
    assert result["tables"][0]["rows"][1] == ["North", "100"]


def test_extract_tables_from_xlsx(tmp_path: Path):
    from openpyxl import Workbook

    path = tmp_path / "workbook.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Plan"
    sheet.append(["Task", "Owner"])
    sheet.append(["Launch", "Ada"])
    workbook.save(path)

    result = svc.extract_tables(path)

    assert result["tables"][0]["caption"] == "Plan"
    assert result["tables"][0]["headers"] == ["Task", "Owner"]
    assert result["tables"][0]["rows"][1] == ["Launch", "Ada"]


def test_ask_with_citations_uses_extractive_fallback(tmp_path: Path):
    path = tmp_path / "policy.md"
    path.write_text(
        "Payment terms\n\nInvoices must be paid within 30 days after receipt.\n\nTermination requires notice.",
        encoding="utf-8",
    )

    result = svc.ask_with_citations(path, "When are invoices paid?", provider_resolver=lambda task="subagent": None)

    assert result["note"] == "extractive_fallback"
    assert "30 days" in result["answer"]
    assert result["citations"]
    assert result["source_blocks"][0]["citation"].startswith("[p1:b")


def test_ask_with_citations_uses_provider_when_available(tmp_path: Path):
    path = tmp_path / "policy.txt"
    path.write_text("Support responses must include citations for source-backed claims.", encoding="utf-8")
    provider = _StubProvider("Use citations for source-backed claims. [p1:b1]")

    result = svc.ask_with_citations(
        path,
        "What must support responses include?",
        provider_resolver=lambda task="subagent": provider,
    )

    assert result["note"] == "llm_qa"
    assert "[p1:b1]" in result["answer"]
    assert provider.calls
    assert "Source blocks" in provider.calls[0][-1]["content"]


def test_compare_documents_reports_added_removed_changed(tmp_path: Path):
    old = tmp_path / "old.txt"
    new = tmp_path / "new.txt"
    old.write_text("Intro\n\nKeep this paragraph.\n\nRemove this paragraph.", encoding="utf-8")
    new.write_text("Intro updated\n\nKeep this paragraph.\n\nAdd this paragraph.", encoding="utf-8")

    result = svc.compare_documents(old, new)

    assert result["changed"][0]["from"]["text"] == "Intro"
    assert result["changed"][0]["to"]["text"] == "Intro updated"
    assert any(item["text"] == "Remove this paragraph." for item in result["removed"])
    assert any(item["text"] == "Add this paragraph." for item in result["added"])


def test_redact_preview_returns_preview_without_writing_file(tmp_path: Path):
    path = tmp_path / "contacts.txt"
    path.write_text("Email ada@example.com or call 212-555-0199.", encoding="utf-8")
    before = path.read_text(encoding="utf-8")

    result = svc.redact_preview(path)

    assert "[REDACTED:EMAIL]" in result["redacted_text"]
    assert "[REDACTED:PHONE]" in result["redacted_text"]
    assert path.read_text(encoding="utf-8") == before


def test_generate_cited_report_fallback_contains_citations(tmp_path: Path):
    path = tmp_path / "brief.txt"
    path.write_text("Market demand increased in the enterprise segment.\n\nCosts remained flat.", encoding="utf-8")

    result = svc.generate_cited_report(path, title="Market Brief", provider_resolver=lambda task="subagent": None)

    assert result["note"] == "extractive_fallback"
    assert result["report"].startswith("# Market Brief")
    assert "[p1:b1]" in result["report"]


def test_edit_docx_dry_run_counts_matches_without_writing(tmp_path: Path):
    path = tmp_path / "memo.docx"
    from docx import Document

    doc = Document()
    doc.add_paragraph("Quarterly memo for ada@example.com")
    doc.save(path)

    preview = svc.edit_docx(path, find="Quarterly", replace="Annual", dry_run=True)
    after = Document(str(path)).paragraphs[0].text

    assert preview["dry_run"] is True
    assert preview["match_count"] == 1
    assert after == "Quarterly memo for ada@example.com"


def test_edit_docx_writes_replacement_when_dry_run_false(tmp_path: Path):
    path = tmp_path / "memo.docx"
    from docx import Document

    doc = Document()
    doc.add_paragraph("Quarterly memo")
    doc.save(path)

    result = svc.edit_docx(path, find="Quarterly", replace="Annual", dry_run=False)
    text = Document(str(path)).paragraphs[0].text

    assert result["ok"] is True
    assert text == "Annual memo"
    assert result["rollback_info"]["backup"]
    assert Path(result["rollback_info"]["backup"]).exists()


def test_edit_docx_replaces_heading_and_table_text(tmp_path: Path):
    from docx import Document

    path = tmp_path / "report.docx"
    doc = Document()
    doc.add_paragraph("Old Title", style="Heading 1")
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Old Title in table"
    doc.save(path)

    preview = svc.edit_docx(path, find="Old Title", replace="New Title", dry_run=True)
    assert preview["match_count"] == 2

    result = svc.edit_docx(path, find="Old Title", replace="New Title", dry_run=False)
    updated = Document(str(path))

    assert result["ok"] is True
    assert updated.paragraphs[0].text == "New Title"
    assert updated.tables[0].cell(0, 0).text == "New Title in table"


def test_edit_docx_dry_run_includes_resource_state(tmp_path: Path):
    path = tmp_path / "memo.docx"
    from docx import Document

    doc = Document()
    doc.add_paragraph("Quarterly memo")
    doc.save(path)

    preview = svc.edit_docx(path, find="Quarterly", replace="Annual", dry_run=True)

    assert preview["_resource_state"]
    assert preview["_resource_state"][0]["path"] == str(path.resolve())


def test_edit_xlsx_dry_run_and_write_cell(tmp_path: Path):
    path = tmp_path / "sheet.xlsx"
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.active.title = "Summary"
    workbook.active["B2"] = "old"
    workbook.save(path)

    preview = svc.edit_xlsx(path, sheet="Summary", cell="B2", value="new", dry_run=True)
    assert preview["dry_run"] is True
    assert preview["diff_preview"][0]["to"] == "new"

    result = svc.edit_xlsx(path, sheet="Summary", cell="B2", value="new", dry_run=False)
    from openpyxl import load_workbook

    assert result["ok"] is True
    assert load_workbook(path).active["B2"].value == "new"
    assert result["rollback_info"]["backup"]
    assert Path(result["rollback_info"]["backup"]).exists()


def test_edit_xlsx_dry_run_includes_resource_state(tmp_path: Path):
    path = tmp_path / "sheet.xlsx"
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.active.title = "Summary"
    workbook.active["B2"] = "old"
    workbook.save(path)

    preview = svc.edit_xlsx(path, sheet="Summary", cell="B2", value="new", dry_run=True)

    assert preview["_resource_state"]
    assert preview["_resource_state"][0]["path"] == str(path.resolve())


def test_document_edit_tools_support_rollback_restore(tmp_path: Path):
    from docx import Document

    from app.core.schemas import ToolResult
    from app.tools import rollback_tools

    path = tmp_path / "memo.docx"
    doc = Document()
    doc.add_paragraph("Original title")
    doc.save(path)

    result = svc.edit_docx(path, find="Original", replace="Changed", dry_run=False)
    assert Document(str(path)).paragraphs[0].text == "Changed title"

    outcome = rollback_tools.rollback_tool_result(
        ToolResult(tool_call_id="doc-edit", ok=True, rollback_info=result["rollback_info"]),
        {"allowed_directories": [str(tmp_path)]},
    )
    assert outcome["ok"] is True
    assert Document(str(path)).paragraphs[0].text == "Original title"


def test_apply_redaction_preview_then_write(tmp_path: Path):
    path = tmp_path / "contacts.txt"
    path.write_text("Email ada@example.com", encoding="utf-8")

    preview = svc.apply_redaction(path, dry_run=True)
    assert preview["dry_run"] is True
    assert preview["preview"]["findings"]

    result = svc.apply_redaction(path, dry_run=False)
    text = path.read_text(encoding="utf-8")

    assert result["ok"] is True
    assert "[REDACTED:EMAIL]" in text
    assert result["rollback_info"]["backup"]


def test_edit_docx_preserves_run_formatting_when_match_is_single_run(tmp_path: Path):
    from docx import Document

    path = tmp_path / "styled.docx"
    doc = Document()
    paragraph = doc.add_paragraph()
    bold_run = paragraph.add_run("Title")
    bold_run.bold = True
    paragraph.add_run(" body")
    doc.save(path)

    result = svc.edit_docx(path, find="Title", replace="Heading", dry_run=False)
    updated = Document(str(path))
    runs = updated.paragraphs[0].runs

    assert result["ok"] is True
    assert runs[0].text == "Heading"
    assert runs[0].bold is True
    assert runs[1].text == " body"


def test_edit_pptx_dry_run_and_write_slide_text(tmp_path: Path):
    from pptx import Presentation

    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Review"
    prs.save(path)

    preview = svc.edit_pptx(path, find="Quarterly", replace="Annual", dry_run=True)
    assert preview["dry_run"] is True
    assert preview["match_count"] == 1
    assert Presentation(str(path)).slides[0].shapes.title.text == "Quarterly Review"

    result = svc.edit_pptx(path, find="Quarterly", replace="Annual", dry_run=False)
    assert result["ok"] is True
    assert Presentation(str(path)).slides[0].shapes.title.text == "Annual Review"
    assert result["rollback_info"]["backup"]


def test_document_intelligence_tools_are_registered_readonly():
    registry = ToolRegistry()
    document_tools.register(registry)

    for name in {
        "document.parse_advanced",
        "document.extract_tables",
        "document.ask_with_citations",
        "document.compare",
        "document.redact_preview",
        "document.generate_cited_report",
    }:
        tool = registry.get(name)
        assert tool.risk_level.value == "R0_READ_ONLY"
        assert tool.is_read_only() is True
        assert tool.requires_authorized_path is True
        assert tool.fast_path_eligible is True

    for name in {
        "document.edit_docx",
        "document.edit_xlsx",
        "document.edit_pptx",
        "document.apply_redaction",
    }:
        tool = registry.get(name)
        assert tool.risk_level.value == "R2_REVERSIBLE_MODIFY"
        assert tool.supports_dry_run is True
        assert tool.is_read_only() is False
        assert tool.requires_authorized_path is True


def test_document_tool_parse_advanced_authorizes_path(tmp_path: Path):
    path = tmp_path / "memo.txt"
    path.write_text("Authorized document text", encoding="utf-8")
    context = {"allowed_directories": [str(tmp_path)]}

    result = document_tools.parse_advanced({"path": str(path)}, context)

    assert result["kind"] == "text"
    assert result["blocks"][0]["text"] == "Authorized document text"
