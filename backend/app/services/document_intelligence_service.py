from __future__ import annotations

import csv
import difflib
import hashlib
import json
import logging
import re
from collections.abc import Callable
from datetime import datetime
from pathlib import Path
from typing import Any
from zipfile import BadZipFile

from app.config import get_env
from app.indexer.ocr_service import IMAGE_EXTENSIONS, extract_pdf_text_with_ocr_fallback, ocr_image_result
from app.observability.best_effort import log_best_effort_failure
from app.orchestration.resource_state import resource_states
from app.policy.redaction import redact_public_text, redact_value
from app.services import document_intelligence_qa as _qa_helpers
from app.services import document_intelligence_text as _text_helpers
from app.services.document_intelligence_models import DocumentBlock, DocumentIR, DocumentTable, ProviderResolver
from app.tools.filesystem_safety import (
    ensure_mutation_path_safe,
    path_exists_or_reparse_point,
    prepare_parent_for_mutation,
    safe_write_text,
)
from app.tools.managed_backups import create_managed_backup
from app.tools.tool_abort import raise_if_tool_aborted

TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".log", ".rst"}
CODE_OR_DATA_EXTENSIONS = {".json", ".yaml", ".yml", ".py", ".ts", ".tsx", ".js"}
TABLE_EXTENSIONS = {".csv", ".xlsx"}
OFFICE_EXTENSIONS = {".pdf", ".docx", ".pptx", ".html", ".htm"}
SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS | CODE_OR_DATA_EXTENSIONS | TABLE_EXTENSIONS | OFFICE_EXTENSIONS | IMAGE_EXTENSIONS
)

DEFAULT_TOP_K = 4
DEFAULT_REPORT_BLOCKS = 8
DEFAULT_PREVIEW_CHARS = 20000
DEFAULT_MAX_PARSE_BYTES = 100 * 1024 * 1024
logger = logging.getLogger(__name__)
_OPTIONAL_IMPORT_ERRORS = (ImportError, OSError)
_TEXT_TABLE_READ_ERRORS = (OSError, UnicodeError, csv.Error)
_OFFICE_LOAD_ERRORS = (OSError, ValueError, KeyError, EOFError, BadZipFile)
_OPTIONAL_HASH_ERRORS = (ImportError, OSError, ValueError, RuntimeError)


class DocumentTooLargeError(ValueError):
    """Raised when a document exceeds the parse size budget."""


class AdvancedParserUnavailable(RuntimeError):
    """Raised when an optional advanced parser cannot be used."""


_blocks_are_changed_pair = _text_helpers._blocks_are_changed_pair
_blocks_from_pages = _text_helpers._blocks_from_pages
_coerce_docling_tables = _text_helpers._coerce_docling_tables
_guess_block_kind = _text_helpers._guess_block_kind
_headers_from_rows = _text_helpers._headers_from_rows
_meaningful_diff_tokens = _text_helpers._meaningful_diff_tokens
_normalize_for_diff = _text_helpers._normalize_for_diff
_normalize_pages = _text_helpers._normalize_pages
_paragraph_blocks = _text_helpers._paragraph_blocks
_rows_from_plain_table_text = _text_helpers._rows_from_plain_table_text
_rows_to_text = _text_helpers._rows_to_text
_split_text_blocks = _text_helpers._split_text_blocks
_stringify_cell = _text_helpers._stringify_cell
_tables_to_text = _text_helpers._tables_to_text

_call_chat = _qa_helpers._call_chat
_document_qa_messages = _qa_helpers._document_qa_messages
_document_report_messages = _qa_helpers._document_report_messages
_fallback_cited_answer = _qa_helpers._fallback_cited_answer
_fallback_cited_report = _qa_helpers._fallback_cited_report
_format_cited_blocks = _qa_helpers._format_cited_blocks
_rank_blocks = _qa_helpers._rank_blocks
_source_block_payload = _qa_helpers._source_block_payload

__all__ = [
    "AdvancedParserUnavailable",
    "DEFAULT_PREVIEW_CHARS",
    "DEFAULT_REPORT_BLOCKS",
    "DEFAULT_TOP_K",
    "DocumentBlock",
    "DocumentIR",
    "DocumentTable",
    "DocumentTooLargeError",
    "ProviderResolver",
    "answer_ir_with_citations",
    "apply_redaction",
    "ask_with_citations",
    "compare_documents",
    "edit_docx",
    "edit_pptx",
    "edit_xlsx",
    "extract_tables",
    "generate_cited_report",
    "parse_advanced",
    "redact_preview",
]


def _append_extraction_warning(
    warnings: list[str],
    message: str,
    operation: str,
    exc: BaseException,
    **context: Any,
) -> None:
    warnings.append(f"{message}: {_redacted_error(exc)}")
    log_best_effort_failure(logger, operation, exc, **context)


def _redacted_error(exc: BaseException) -> str:
    return redact_public_text(str(redact_value(str(exc)) or ""))


def _max_parse_bytes() -> int:
    raw = get_env("LENGRVIS_DOCUMENT_MAX_PARSE_BYTES")
    if not raw:
        return DEFAULT_MAX_PARSE_BYTES
    try:
        return max(1, int(raw))
    except ValueError:
        return DEFAULT_MAX_PARSE_BYTES


def _ensure_parseable_file_size(document_path: Path) -> None:
    try:
        size = document_path.stat().st_size
    except OSError as exc:
        raise FileNotFoundError(f"Document not found: {document_path}") from exc
    max_bytes = _max_parse_bytes()
    if size > max_bytes:
        raise DocumentTooLargeError(f"Document exceeds parse size limit ({size} bytes; max {max_bytes}).")


def parse_advanced(
    path: str | Path,
    *,
    settings: Any | None = None,
    try_advanced: bool = True,
) -> DocumentIR:
    document_path = Path(path)
    if not document_path.exists():
        raise FileNotFoundError(f"Document not found: {document_path}")
    _ensure_parseable_file_size(document_path)

    warnings: list[str] = []
    if try_advanced:
        for parser_name, parser in (("docling", _parse_with_docling), ("unstructured", _parse_with_unstructured)):
            try:
                ir = parser(document_path)
            except AdvancedParserUnavailable as exc:
                warnings.append(f"{parser_name} parser unavailable: {_redacted_error(exc)}")
                continue
            except Exception as exc:  # noqa: BLE001 - broad-exception-boundary: parser plugins fail in many environment-specific ways.
                _append_extraction_warning(
                    warnings,
                    f"{parser_name} parser failed",
                    "document_intelligence.advanced_parser",
                    exc,
                    parser=parser_name,
                    extension=document_path.suffix.lower(),
                )
                continue
            if ir.blocks or ir.tables:
                ir.warnings = warnings + ir.warnings
                return ir
            warnings.append(f"{parser_name} parser produced no content.")

    fallback = _parse_with_builtin(document_path, settings=settings, warnings=warnings)
    return fallback


def extract_tables(path: str | Path, *, settings: Any | None = None) -> dict[str, Any]:
    ir = parse_advanced(path, settings=settings)
    return {
        "document_id": ir.document_id,
        "path": ir.path,
        "kind": ir.kind,
        "tables": [table.as_dict() for table in ir.tables],
        "parse_engine": ir.parse_engine,
        "warnings": ir.warnings,
    }


def ask_with_citations(
    path: str | Path,
    question: str,
    *,
    settings: Any | None = None,
    provider_resolver: ProviderResolver | None = None,
    top_k: int = DEFAULT_TOP_K,
) -> dict[str, Any]:
    ir = parse_advanced(path, settings=settings)
    return answer_ir_with_citations(
        ir,
        question,
        provider_resolver=provider_resolver,
        top_k=top_k,
    )


def answer_ir_with_citations(
    ir: DocumentIR,
    question: str,
    *,
    provider_resolver: ProviderResolver | None = None,
    top_k: int = DEFAULT_TOP_K,
) -> dict[str, Any]:
    cleaned_question = (question or "").strip()
    if not cleaned_question:
        return {
            "document_id": ir.document_id,
            "question": cleaned_question,
            "answer": "",
            "note": "no_question",
            "citations": [],
            "source_blocks": [],
            "warnings": ir.warnings,
        }

    ranked = _rank_blocks(cleaned_question, ir.blocks, top_k=top_k)
    if not ranked:
        return {
            "document_id": ir.document_id,
            "question": cleaned_question,
            "answer": "",
            "note": "no_relevant_blocks",
            "citations": [],
            "source_blocks": [],
            "warnings": ir.warnings,
        }

    prompt_context = _format_cited_blocks(ranked, max_chars=9000)
    messages = _document_qa_messages(cleaned_question, prompt_context)
    answer = _call_chat(messages, provider_resolver=provider_resolver)
    if not answer:
        answer = _fallback_cited_answer(cleaned_question, ranked)
        note = "extractive_fallback"
    else:
        note = "llm_qa"

    return {
        "document_id": ir.document_id,
        "question": cleaned_question,
        "answer": answer,
        "note": note,
        "citations": [block.citation for block in ranked],
        "source_blocks": [_source_block_payload(block) for block in ranked],
        "warnings": ir.warnings,
    }


def compare_documents(
    left_path: str | Path,
    right_path: str | Path,
    *,
    settings: Any | None = None,
) -> dict[str, Any]:
    left = parse_advanced(left_path, settings=settings)
    right = parse_advanced(right_path, settings=settings)
    left_blocks = _paragraph_blocks(left)
    right_blocks = _paragraph_blocks(right)
    matcher = difflib.SequenceMatcher(
        None,
        [_normalize_for_diff(block.text) for block in left_blocks],
        [_normalize_for_diff(block.text) for block in right_blocks],
        autojunk=False,
    )

    added: list[dict[str, Any]] = []
    removed: list[dict[str, Any]] = []
    changed: list[dict[str, Any]] = []
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == "equal":
            continue
        if tag == "delete":
            removed.extend(_source_block_payload(block) for block in left_blocks[i1:i2])
            continue
        if tag == "insert":
            added.extend(_source_block_payload(block) for block in right_blocks[j1:j2])
            continue

        old_slice = left_blocks[i1:i2]
        new_slice = right_blocks[j1:j2]
        paired_count = min(len(old_slice), len(new_slice))
        for index in range(paired_count):
            old_block = old_slice[index]
            new_block = new_slice[index]
            if _blocks_are_changed_pair(old_block.text, new_block.text):
                changed.append(
                    {
                        "from": _source_block_payload(old_block),
                        "to": _source_block_payload(new_block),
                    }
                )
            else:
                removed.append(_source_block_payload(old_block))
                added.append(_source_block_payload(new_block))
        removed.extend(_source_block_payload(block) for block in old_slice[paired_count:])
        added.extend(_source_block_payload(block) for block in new_slice[paired_count:])

    return {
        "left": {"document_id": left.document_id, "path": left.path},
        "right": {"document_id": right.document_id, "path": right.path},
        "added": added,
        "removed": removed,
        "changed": changed,
        "warnings": left.warnings + right.warnings,
    }


def redact_preview(
    path: str | Path,
    *,
    settings: Any | None = None,
    custom_patterns: dict[str, str] | None = None,
    max_chars: int = DEFAULT_PREVIEW_CHARS,
) -> dict[str, Any]:
    ir = parse_advanced(path, settings=settings)
    source = ir.text[: max(1, max_chars)]
    redacted = source
    findings: list[dict[str, Any]] = []

    for label, pattern in _redaction_patterns(custom_patterns).items():
        regex = re.compile(pattern)
        matches = list(regex.finditer(redacted))
        if not matches:
            continue
        redacted = regex.sub(f"[REDACTED:{label}]", redacted)
        findings.append({"type": label, "count": len(matches)})

    return {
        "document_id": ir.document_id,
        "path": ir.path,
        "redacted_text": redacted,
        "truncated": len(ir.text) > max_chars,
        "findings": findings,
        "warnings": ir.warnings,
    }


def generate_cited_report(
    path: str | Path,
    *,
    title: str = "Cited Report",
    query: str = "",
    settings: Any | None = None,
    provider_resolver: ProviderResolver | None = None,
    max_blocks: int = DEFAULT_REPORT_BLOCKS,
) -> dict[str, Any]:
    ir = parse_advanced(path, settings=settings)
    selected = _rank_blocks(query, ir.blocks, top_k=max_blocks) if query.strip() else _paragraph_blocks(ir)[:max_blocks]
    if not selected:
        return {
            "document_id": ir.document_id,
            "title": title,
            "report": f"# {title}\n\nNo source blocks were available.",
            "note": "no_source_blocks",
            "citations": [],
            "warnings": ir.warnings,
        }

    source = _format_cited_blocks(selected, max_chars=14000)
    messages = _document_report_messages(title, source)
    report = _call_chat(messages, provider_resolver=provider_resolver)
    if report:
        note = "llm_report"
    else:
        report = _fallback_cited_report(title, selected)
        note = "extractive_fallback"

    return {
        "document_id": ir.document_id,
        "title": title,
        "report": report,
        "note": note,
        "citations": [block.citation for block in selected],
        "source_blocks": [_source_block_payload(block) for block in selected],
        "warnings": ir.warnings,
    }


def _parse_with_docling(path: Path) -> DocumentIR:
    try:
        from docling.document_converter import DocumentConverter
    except _OPTIONAL_IMPORT_ERRORS as exc:
        raise AdvancedParserUnavailable(str(exc)) from exc

    converter = DocumentConverter()
    result = converter.convert(str(path))
    document = getattr(result, "document", result)
    text = ""
    for method_name in ("export_to_markdown", "export_to_text"):
        method = getattr(document, method_name, None)
        if callable(method):
            text = str(method() or "").strip()
            if text:
                break
    if not text:
        text = str(document or "").strip()
    if not text:
        raise AdvancedParserUnavailable("docling returned an empty document")

    tables = _coerce_docling_tables(getattr(document, "tables", []) or [])
    return _build_ir(
        path,
        kind=_detect_kind(path),
        pages=[{"page": 1, "text": text}],
        tables=tables,
        parse_engine="docling",
        metadata={"advanced_parser": "docling"},
    )


def _parse_with_unstructured(path: Path) -> DocumentIR:
    try:
        from unstructured.partition.auto import partition
    except _OPTIONAL_IMPORT_ERRORS as exc:
        raise AdvancedParserUnavailable(str(exc)) from exc

    elements = partition(filename=str(path))
    if not elements:
        raise AdvancedParserUnavailable("unstructured returned no elements")

    pages: dict[int, list[str]] = {}
    tables: list[DocumentTable] = []
    for element in elements:
        text = str(element or "").strip()
        metadata = getattr(element, "metadata", None)
        page_number = int(getattr(metadata, "page_number", None) or 1)
        if text:
            pages.setdefault(page_number, []).append(text)
        category = str(getattr(element, "category", element.__class__.__name__) or "")
        if category.casefold() == "table" and text:
            rows = _rows_from_plain_table_text(text)
            if rows:
                tables.append(
                    DocumentTable(
                        id=f"table-{len(tables) + 1}",
                        rows=rows,
                        headers=_headers_from_rows(rows),
                        page=page_number,
                        caption="",
                        metadata={"source": "unstructured"},
                    )
                )

    page_payload = [{"page": page, "text": "\n\n".join(parts)} for page, parts in sorted(pages.items())]
    return _build_ir(
        path,
        kind=_detect_kind(path),
        pages=page_payload,
        tables=tables,
        parse_engine="unstructured",
        metadata={"advanced_parser": "unstructured"},
    )


def _parse_with_builtin(path: Path, *, settings: Any | None, warnings: list[str]) -> DocumentIR:
    ext = path.suffix.lower()
    kind = _detect_kind(path)
    ocr_engine = ""
    tables: list[DocumentTable] = []
    metadata: dict[str, Any] = {}

    if ext in TEXT_EXTENSIONS or ext in {".yaml", ".yml", ".py", ".ts", ".tsx", ".js"}:
        text = path.read_text(encoding="utf-8", errors="ignore")
        pages = [{"page": 1, "text": text}]
    elif ext == ".json":
        text, metadata = _read_json_text(path, warnings)
        pages = [{"page": 1, "text": text}]
    elif ext == ".csv":
        tables = _extract_csv_tables(path, warnings)
        text = _tables_to_text(tables)
        pages = [{"page": 1, "text": text}]
    elif ext == ".xlsx":
        pages, tables = _extract_xlsx(path, warnings)
    elif ext == ".pdf":
        pages = _extract_pdf_pages(path, settings=settings, warnings=warnings)
    elif ext == ".docx":
        pages, tables = _extract_docx(path, warnings)
    elif ext == ".pptx":
        pages, tables = _extract_pptx(path, warnings)
    elif ext in {".html", ".htm"}:
        pages, tables = _extract_html(path, warnings)
    elif ext in IMAGE_EXTENSIONS:
        result = ocr_image_result(path, settings=settings)
        ocr_engine = result.source
        if not result.ok:
            warnings.append(result.error or "Image OCR produced no text.")
        pages = [{"page": 1, "text": result.text}]
    else:
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
        except OSError as exc:
            text = ""
            _append_extraction_warning(
                warnings,
                f"Unsupported document type: {ext or '<none>'}",
                "document_intelligence.unsupported_text_fallback",
                exc,
                extension=ext or "<none>",
            )
        pages = [{"page": 1, "text": text}]

    return _build_ir(
        path,
        kind=kind,
        pages=pages,
        tables=tables,
        parse_engine="builtin",
        ocr_engine=ocr_engine,
        metadata=metadata,
        warnings=warnings,
    )


def _extract_csv_tables(path: Path, warnings: list[str]) -> list[DocumentTable]:
    try:
        with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as handle:
            rows = [[_stringify_cell(cell) for cell in row] for row in csv.reader(handle)]
    except _TEXT_TABLE_READ_ERRORS as exc:
        _append_extraction_warning(
            warnings,
            "CSV extraction failed",
            "document_intelligence.csv",
            exc,
            extension=path.suffix.lower(),
        )
        return []
    if not rows:
        return []
    return [
        DocumentTable(
            id="table-1",
            rows=rows,
            headers=_headers_from_rows(rows),
            page=1,
            caption=path.name,
            metadata={"source": "csv"},
        )
    ]


def _extract_xlsx(path: Path, warnings: list[str]) -> tuple[list[dict[str, Any]], list[DocumentTable]]:
    try:
        from openpyxl import load_workbook
        from openpyxl.utils.exceptions import InvalidFileException
    except _OPTIONAL_IMPORT_ERRORS as exc:
        _append_extraction_warning(
            warnings,
            "XLSX extraction unavailable",
            "document_intelligence.xlsx.import",
            exc,
            extension=path.suffix.lower(),
        )
        return [{"page": 1, "text": ""}], []

    pages: list[dict[str, Any]] = []
    tables: list[DocumentTable] = []
    try:
        workbook = load_workbook(path, read_only=True, data_only=True)
    except (*_OFFICE_LOAD_ERRORS, InvalidFileException) as exc:
        _append_extraction_warning(
            warnings,
            "XLSX extraction failed",
            "document_intelligence.xlsx.load",
            exc,
            extension=path.suffix.lower(),
        )
        return [{"page": 1, "text": ""}], []

    try:
        for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
            rows = [
                [_stringify_cell(value) for value in row]
                for row in sheet.iter_rows(values_only=True)
                if any(value is not None and str(value).strip() for value in row)
            ]
            if rows:
                tables.append(
                    DocumentTable(
                        id=f"table-{len(tables) + 1}",
                        rows=rows,
                        headers=_headers_from_rows(rows),
                        page=sheet_index,
                        caption=sheet.title,
                        metadata={"source": "xlsx", "sheet": sheet.title},
                    )
                )
            pages.append(
                {
                    "page": sheet_index,
                    "text": _rows_to_text(rows, title=f"Sheet: {sheet.title}"),
                    "metadata": {"sheet": sheet.title},
                }
            )
    finally:
        close = getattr(workbook, "close", None)
        if callable(close):
            close()

    return pages or [{"page": 1, "text": ""}], tables


def _extract_pdf_pages(path: Path, *, settings: Any | None, warnings: list[str]) -> list[dict[str, Any]]:
    try:
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError
    except _OPTIONAL_IMPORT_ERRORS as exc:
        _append_extraction_warning(
            warnings,
            "PDF page extraction fell back to OCR/text fallback",
            "document_intelligence.pdf.page_extract",
            exc,
            extension=path.suffix.lower(),
        )
    else:
        try:
            reader = PdfReader(str(path))
            pages = [
                {"page": index, "text": page.extract_text() or ""} for index, page in enumerate(reader.pages, start=1)
            ]
            if any(len(page["text"].strip()) >= 24 for page in pages):
                return pages
        except (*_OFFICE_LOAD_ERRORS, RuntimeError, PdfReadError) as exc:
            _append_extraction_warning(
                warnings,
                "PDF page extraction fell back to OCR/text fallback",
                "document_intelligence.pdf.page_extract",
                exc,
                extension=path.suffix.lower(),
            )

    text = extract_pdf_text_with_ocr_fallback(path, settings=settings)
    return [{"page": 1, "text": text}]


def _extract_docx(path: Path, warnings: list[str]) -> tuple[list[dict[str, Any]], list[DocumentTable]]:
    try:
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
    except _OPTIONAL_IMPORT_ERRORS as exc:
        _append_extraction_warning(
            warnings,
            "DOCX extraction unavailable",
            "document_intelligence.docx.import",
            exc,
            extension=path.suffix.lower(),
        )
        return [{"page": 1, "text": ""}], []

    try:
        document = Document(str(path))
    except (*_OFFICE_LOAD_ERRORS, PackageNotFoundError) as exc:
        _append_extraction_warning(
            warnings,
            "DOCX extraction failed",
            "document_intelligence.docx.load",
            exc,
            extension=path.suffix.lower(),
        )
        return [{"page": 1, "text": ""}], []

    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    tables: list[DocumentTable] = []
    for table in document.tables:
        rows = [[_stringify_cell(cell.text) for cell in row.cells] for row in table.rows]
        if rows:
            tables.append(
                DocumentTable(
                    id=f"table-{len(tables) + 1}",
                    rows=rows,
                    headers=_headers_from_rows(rows),
                    page=1,
                    metadata={"source": "docx"},
                )
            )
    table_text = _tables_to_text(tables)
    text = "\n\n".join(part for part in ["\n\n".join(paragraphs), table_text] if part)
    return [{"page": 1, "text": text}], tables


def _extract_pptx(path: Path, warnings: list[str]) -> tuple[list[dict[str, Any]], list[DocumentTable]]:
    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except _OPTIONAL_IMPORT_ERRORS as exc:
        _append_extraction_warning(
            warnings,
            "PPTX extraction unavailable",
            "document_intelligence.pptx.import",
            exc,
            extension=path.suffix.lower(),
        )
        return [{"page": 1, "text": ""}], []

    try:
        presentation = Presentation(str(path))
    except (*_OFFICE_LOAD_ERRORS, PackageNotFoundError) as exc:
        _append_extraction_warning(
            warnings,
            "PPTX extraction failed",
            "document_intelligence.pptx.load",
            exc,
            extension=path.suffix.lower(),
        )
        return [{"page": 1, "text": ""}], []

    pages: list[dict[str, Any]] = []
    tables: list[DocumentTable] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                lines.append(str(text))
            if getattr(shape, "has_table", False):
                rows = [[_stringify_cell(cell.text) for cell in row.cells] for row in shape.table.rows]
                if rows:
                    tables.append(
                        DocumentTable(
                            id=f"table-{len(tables) + 1}",
                            rows=rows,
                            headers=_headers_from_rows(rows),
                            page=slide_index,
                            metadata={"source": "pptx", "slide": slide_index},
                        )
                    )
        pages.append({"page": slide_index, "text": "\n\n".join(lines), "metadata": {"slide": slide_index}})
    return pages or [{"page": 1, "text": ""}], tables


def _extract_html(path: Path, warnings: list[str]) -> tuple[list[dict[str, Any]], list[DocumentTable]]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    try:
        from bs4 import BeautifulSoup
    except _OPTIONAL_IMPORT_ERRORS as exc:
        _append_extraction_warning(
            warnings,
            "HTML parser unavailable, used tag-stripping fallback",
            "document_intelligence.html.import",
            exc,
            extension=path.suffix.lower(),
        )
        text = re.sub(r"<[^>]+>", " ", raw)
        return [{"page": 1, "text": " ".join(text.split())}], []

    soup = BeautifulSoup(raw, "html.parser")
    tables: list[DocumentTable] = []
    for table in soup.find_all("table"):
        rows: list[list[str]] = []
        for tr in table.find_all("tr"):
            cells = tr.find_all(["th", "td"])
            row = [" ".join(cell.get_text(" ", strip=True).split()) for cell in cells]
            if row:
                rows.append(row)
        if rows:
            tables.append(
                DocumentTable(
                    id=f"table-{len(tables) + 1}",
                    rows=rows,
                    headers=_headers_from_rows(rows),
                    page=1,
                    metadata={"source": "html"},
                )
            )
    text = soup.get_text("\n", strip=True)
    return [{"page": 1, "text": text}], tables


def _build_ir(
    path: Path,
    *,
    kind: str,
    pages: list[dict[str, Any]],
    tables: list[DocumentTable],
    parse_engine: str,
    ocr_engine: str = "",
    metadata: dict[str, Any] | None = None,
    warnings: list[str] | None = None,
) -> DocumentIR:
    warning_list = list(warnings or [])
    document_id = _document_id(path, warning_list)
    normalized_pages = _normalize_pages(pages)
    blocks = _blocks_from_pages(normalized_pages)
    file_metadata = _file_metadata(path)
    file_metadata.update(metadata or {})
    file_metadata["page_count"] = len(normalized_pages)
    file_metadata["block_count"] = len(blocks)
    file_metadata["table_count"] = len(tables)
    return DocumentIR(
        document_id=document_id,
        path=str(path),
        kind=kind,
        pages=normalized_pages,
        blocks=blocks,
        tables=tables,
        metadata=file_metadata,
        parse_engine=parse_engine,
        ocr_engine=ocr_engine,
        warnings=warning_list,
    )


def _document_id(path: Path, warnings: list[str]) -> str:
    data = path.read_bytes()
    try:
        import blake3

        return f"blake3:{blake3.blake3(data).hexdigest()}"
    except _OPTIONAL_HASH_ERRORS as exc:
        _append_extraction_warning(
            warnings,
            "blake3 unavailable, used sha256 fallback",
            "document_intelligence.document_id.blake3",
            exc,
            extension=path.suffix.lower(),
        )
        return f"sha256:{hashlib.sha256(data).hexdigest()}"


def _file_metadata(path: Path) -> dict[str, Any]:
    stat = path.stat()
    return {
        "name": path.name,
        "extension": path.suffix.lower(),
        "size_bytes": stat.st_size,
        "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds"),
    }


def _detect_kind(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in {".md", ".markdown"}:
        return "markdown"
    if ext in TEXT_EXTENSIONS:
        return "text"
    if ext == ".json":
        return "json"
    if ext == ".csv":
        return "csv"
    if ext == ".xlsx":
        return "xlsx"
    if ext == ".pdf":
        return "pdf"
    if ext == ".docx":
        return "docx"
    if ext == ".pptx":
        return "pptx"
    if ext in {".html", ".htm"}:
        return "html"
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in CODE_OR_DATA_EXTENSIONS:
        return "text"
    return ext.lstrip(".") or "unknown"


def _read_json_text(path: Path, warnings: list[str]) -> tuple[str, dict[str, Any]]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError as exc:
        warnings.append(f"JSON parse failed, returned raw text: {exc}")
        return raw, {}
    return json.dumps(parsed, ensure_ascii=False, indent=2), {"json_type": type(parsed).__name__}


def _backup_document(path: Path, abort_context: dict[str, Any] | None = None) -> dict[str, str | int | bool] | None:
    if not path.is_file():
        return None
    _ensure_document_mutation_safe(path, abort_context)
    raise_if_tool_aborted(abort_context)
    return create_managed_backup(path)


def _document_allowed_directories(context: dict[str, Any] | None, path: Path | None = None) -> list[str]:
    allowed = [str(path) for path in (context or {}).get("allowed_directories") or []]
    if allowed or path is None:
        return allowed
    return [str(path.parent)]


def _ensure_document_mutation_safe(path: Path, context: dict[str, Any] | None) -> None:
    allowed = _document_allowed_directories(context, path)
    ensure_mutation_path_safe(path, allowed, include_self=True, context=context)


def _prepare_document_save(path: Path, context: dict[str, Any] | None) -> None:
    allowed = _document_allowed_directories(context, path)
    prepare_parent_for_mutation(path, allowed, context)
    ensure_mutation_path_safe(path, allowed, include_self=path_exists_or_reparse_point(path), context=context)


def _safe_save_office_document(path: Path, save: Callable[[], None], context: dict[str, Any] | None) -> None:
    _prepare_document_save(path, context)
    raise_if_tool_aborted(context)
    save()
    ensure_mutation_path_safe(path, _document_allowed_directories(context, path), include_self=True, context=context)


def _document_resource_state(path: Path) -> list[dict[str, Any]]:
    return resource_states([path])


def _replace_in_paragraph_preserve_runs(paragraph: Any, needle: str, replacement: str, *, dry_run: bool) -> int:
    if needle not in paragraph.text:
        return 0
    count = paragraph.text.count(needle)
    if dry_run:
        return count
    runs = list(paragraph.runs)
    for run in runs:
        if needle in run.text:
            run.text = run.text.replace(needle, replacement)
    if needle not in paragraph.text:
        return count
    new_text = paragraph.text.replace(needle, replacement)
    if not runs:
        paragraph.add_run(new_text)
        return count
    runs[0].text = new_text
    for run in runs[1:]:
        run.text = ""
    return count


def _docx_replace_in_paragraph(paragraph: Any, needle: str, replacement: str, *, dry_run: bool) -> int:
    return _replace_in_paragraph_preserve_runs(paragraph, needle, replacement, dry_run=dry_run)


def _docx_replace_text(doc: Any, needle: str, replacement: str, *, dry_run: bool) -> int:
    match_count = 0
    for paragraph in doc.paragraphs:
        match_count += _docx_replace_in_paragraph(paragraph, needle, replacement, dry_run=dry_run)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    match_count += _docx_replace_in_paragraph(paragraph, needle, replacement, dry_run=dry_run)
    return match_count


def _pptx_replace_text(prs: Any, needle: str, replacement: str, *, dry_run: bool) -> int:
    match_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                match_count += _replace_in_paragraph_preserve_runs(paragraph, needle, replacement, dry_run=dry_run)
    return match_count


def edit_pptx(
    path: str | Path,
    *,
    find: str,
    replace: str,
    dry_run: bool = True,
    abort_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    from pptx import Presentation

    path_obj = Path(path)
    needle = str(find or "")
    if not needle:
        raise ValueError("document.edit_pptx requires a non-empty 'find' string.")
    _ensure_parseable_file_size(path_obj)
    replacement = str(replace or "")
    prs = Presentation(str(path_obj))
    match_count = _pptx_replace_text(prs, needle, replacement, dry_run=True)
    diff_preview = [
        {
            "action": "edit_pptx",
            "path": str(path_obj),
            "find": needle,
            "replace": replacement,
            "match_count": match_count,
        }
    ]
    if dry_run:
        return {
            "dry_run": True,
            "path": str(path_obj),
            "match_count": match_count,
            "diff_preview": diff_preview,
            "_resource_state": _document_resource_state(path_obj),
        }
    if match_count <= 0:
        return {"ok": False, "error_code": "NO_MATCH", "path": str(path_obj), "match_count": 0}
    raise_if_tool_aborted(abort_context)
    backup = _backup_document(path_obj, abort_context)
    prs = Presentation(str(path_obj))
    _pptx_replace_text(prs, needle, replacement, dry_run=False)
    _safe_save_office_document(path_obj, lambda: prs.save(str(path_obj)), abort_context)
    return {
        "ok": True,
        "path": str(path_obj),
        "match_count": match_count,
        "changed_paths": [str(path_obj)],
        "diff_preview": diff_preview,
        "rollback_info": {"backup": backup},
    }


def apply_redaction(
    path: str | Path,
    *,
    settings: Any | None = None,
    custom_patterns: dict[str, str] | None = None,
    max_chars: int = DEFAULT_PREVIEW_CHARS,
    dry_run: bool = True,
    abort_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    path_obj = Path(path)
    preview = redact_preview(
        path_obj,
        settings=settings,
        custom_patterns=custom_patterns,
        max_chars=max_chars,
    )
    diff_preview = [
        {
            "action": "apply_redaction",
            "path": str(path_obj),
            "findings": preview.get("findings") or [],
        }
    ]
    if dry_run:
        return {
            "dry_run": True,
            "path": str(path_obj),
            "preview": preview,
            "diff_preview": diff_preview,
            "findings": preview.get("findings") or [],
            "_resource_state": _document_resource_state(path_obj),
        }
    raise_if_tool_aborted(abort_context)
    backup = _backup_document(path_obj, abort_context)
    ext = path_obj.suffix.lower()
    if ext in {".txt", ".md", ".csv", ".json"}:
        safe_write_text(
            path_obj,
            str(preview.get("redacted_text") or ""),
            _document_allowed_directories(abort_context, path_obj),
            abort_context,
        )
    elif ext == ".docx":
        _apply_redaction_docx(path_obj, custom_patterns, abort_context=abort_context)
    else:
        raise ValueError(f"document.apply_redaction does not support {ext} files yet.")
    return {
        "ok": True,
        "path": str(path_obj),
        "findings": preview.get("findings") or [],
        "changed_paths": [str(path_obj)],
        "diff_preview": diff_preview,
        "rollback_info": {"backup": backup},
    }


def edit_docx(
    path: str | Path,
    *,
    find: str,
    replace: str,
    dry_run: bool = True,
    abort_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    from docx import Document

    path_obj = Path(path)
    needle = str(find or "")
    if not needle:
        raise ValueError("document.edit_docx requires a non-empty 'find' string.")
    _ensure_parseable_file_size(path_obj)
    replacement = str(replace or "")
    doc = Document(str(path_obj))
    match_count = _docx_replace_text(doc, needle, replacement, dry_run=True)
    diff_preview = [
        {
            "action": "edit_docx",
            "path": str(path_obj),
            "find": needle,
            "replace": replacement,
            "match_count": match_count,
        }
    ]
    if dry_run:
        return {
            "dry_run": True,
            "path": str(path_obj),
            "match_count": match_count,
            "diff_preview": diff_preview,
            "_resource_state": _document_resource_state(path_obj),
        }
    if match_count <= 0:
        return {"ok": False, "error_code": "NO_MATCH", "path": str(path_obj), "match_count": 0}
    raise_if_tool_aborted(abort_context)
    backup = _backup_document(path_obj, abort_context)
    doc = Document(str(path_obj))
    _docx_replace_text(doc, needle, replacement, dry_run=False)
    _safe_save_office_document(path_obj, lambda: doc.save(str(path_obj)), abort_context)
    return {
        "ok": True,
        "path": str(path_obj),
        "match_count": match_count,
        "changed_paths": [str(path_obj)],
        "diff_preview": diff_preview,
        "rollback_info": {"backup": backup},
    }


def edit_xlsx(
    path: str | Path,
    *,
    sheet: str,
    cell: str,
    value: Any,
    dry_run: bool = True,
    abort_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    from openpyxl import load_workbook

    path_obj = Path(path)
    sheet_name = str(sheet or "").strip()
    cell_ref = str(cell or "").strip().upper()
    if not sheet_name or not cell_ref:
        raise ValueError("document.edit_xlsx requires non-empty 'sheet' and 'cell'.")
    _ensure_parseable_file_size(path_obj)
    workbook = load_workbook(path_obj)
    if sheet_name not in workbook.sheetnames:
        raise ValueError(f"Sheet not found: {sheet_name}")
    worksheet = workbook[sheet_name]
    previous = worksheet[cell_ref].value
    diff_preview = [
        {
            "action": "edit_xlsx",
            "path": str(path_obj),
            "sheet": sheet_name,
            "cell": cell_ref,
            "from": previous,
            "to": value,
        }
    ]
    if dry_run:
        return {
            "dry_run": True,
            "path": str(path_obj),
            "sheet": sheet_name,
            "cell": cell_ref,
            "diff_preview": diff_preview,
            "_resource_state": _document_resource_state(path_obj),
        }
    raise_if_tool_aborted(abort_context)
    backup = _backup_document(path_obj, abort_context)
    worksheet[cell_ref].value = value
    _safe_save_office_document(path_obj, lambda: workbook.save(path_obj), abort_context)
    return {
        "ok": True,
        "path": str(path_obj),
        "sheet": sheet_name,
        "cell": cell_ref,
        "changed_paths": [str(path_obj)],
        "diff_preview": diff_preview,
        "rollback_info": {"backup": backup},
    }


def _apply_redaction_docx(
    path: Path,
    custom_patterns: dict[str, str] | None,
    *,
    abort_context: dict[str, Any] | None = None,
) -> None:
    from docx import Document

    doc = Document(str(path))
    patterns = _redaction_patterns(custom_patterns)
    for paragraph in doc.paragraphs:
        text = paragraph.text
        for label, pattern in patterns.items():
            text = re.compile(pattern).sub(f"[REDACTED:{label}]", text)
        paragraph.text = text
    _safe_save_office_document(path, lambda: doc.save(str(path)), abort_context)


def _redaction_patterns(custom_patterns: dict[str, str] | None) -> dict[str, str]:
    patterns = {
        "EMAIL": r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b",
        "SSN": r"\b\d{3}-\d{2}-\d{4}\b",
        "PHONE": r"(?<!\d)(?:\+?\d{1,3}[-.\s]?)?(?:\(?\d{3}\)?[-.\s]?)\d{3}[-.\s]?\d{4}(?!\d)",
        "CREDIT_CARD": r"\b(?:\d[ -]*?){13,16}\b",
        "CN_ID": r"\b\d{17}[\dXx]\b",
    }
    patterns.update(custom_patterns or {})
    return patterns
