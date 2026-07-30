from __future__ import annotations

import csv
from pathlib import Path
from typing import Any

from app.commerce.entitlements import Feature, active_plan, has_feature, require_feature
from app.commerce.licensing import commercial_release_enabled
from app.config import get_env
from app.core.content_provenance import record_tool_output_provenance
from app.core.paths import resolve_authorized
from app.indexer.ocr_service import extract_pdf_text_with_ocr_fallback
from app.llm.registry import LOCAL_PROVIDERS, get_effective_settings
from app.policy.privacy import can_upload_file_content
from app.policy.risk import RiskLevel
from app.services import document_intelligence_service, document_service
from app.tools.schemas import ToolDefinition
from app.tools.tool_abort import raise_if_tool_aborted
from app.tools.tool_catalog import tool_description, tool_search_hint

_EXTRACT_TEXT_LIMIT = 20000
_CHUNK_CHARS = document_service.DEFAULT_CHUNK_CHARS
DEFAULT_MAX_CSV_ANALYSIS_ROWS = 100_000


def _allowed(context: dict[str, Any]) -> list[str]:
    return list(context.get("allowed_directories") or [])


def _document_disabled() -> dict[str, Any]:
    return {"error": "Document AI is not available on this plan."}


def _require_document_ai_tool(fn):
    def execute(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
        settings = context.get("settings") or get_effective_settings()
        plan = active_plan(settings)
        if commercial_release_enabled():
            require_feature(plan, Feature.DOCUMENT_AI)
        elif not has_feature(plan, Feature.DOCUMENT_AI):
            return _document_disabled()
        return fn(args, context)

    return execute


def _document_max_chars_to_llm(context: dict[str, Any]) -> int:
    settings = context.get("settings")
    value = getattr(settings, "document_max_chars_to_llm", document_service.DEFAULT_MAX_CHARS_TO_LLM)
    return max(1, int(value))


def _max_csv_analysis_rows() -> int:
    raw = get_env("LENGRVIS_DOCUMENT_MAX_CSV_ROWS")
    if raw is None:
        return DEFAULT_MAX_CSV_ANALYSIS_ROWS
    try:
        return max(1, int(raw))
    except ValueError:
        return DEFAULT_MAX_CSV_ANALYSIS_ROWS


def _provider(task: str = "subagent"):
    return document_service._provider(task)


def _document_provider(context: dict[str, Any]):
    settings = context.get("settings")
    if settings is None:
        return _provider

    decision = can_upload_file_content(settings)
    if decision.allowed or _document_provider_is_local(settings):
        return _provider

    def resolver(task: str = "subagent"):
        provider = _provider(task)
        if _provider_is_cloud(provider):
            return None
        return provider

    return resolver


def _document_provider_is_local(settings: Any) -> bool:
    mode = str(getattr(settings, "mode", "efficiency") or "efficiency").lower()
    provider_name = str(getattr(settings, "provider_name", "") or "").lower()
    return mode in {"privacy"} or provider_name in LOCAL_PROVIDERS


def _provider_is_cloud(provider: Any) -> bool:
    name = str(getattr(provider, "name", "") or provider.__class__.__name__).casefold()
    module = str(getattr(provider.__class__, "__module__", "") or "").casefold()
    markers = ("openai", "azure", "deepseek", "hunyuan", "custom_http")
    return any(marker in name or marker in module for marker in markers)


def _chunk_text(text: str, chunk_chars: int = _CHUNK_CHARS) -> list[str]:
    return document_service.chunk_document(
        text,
        chunk_chars=chunk_chars,
        overlap=0,
        max_chunks=None,
        max_chars=max(len(text or ""), 1),
    )


def _rank_chunks(query: str, chunks: list[str]) -> list[str]:
    return [chunk.text for chunk in document_service.rank_chunks(query, chunks, top_k=len(chunks) or 1)]


def extract_text_from_path(path: Path) -> str:
    document_intelligence_service._ensure_parseable_file_size(path)
    ext = path.suffix.lower()
    if ext in {".txt", ".md", ".json", ".csv", ".py", ".ts", ".tsx", ".js", ".yaml", ".yml"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if ext == ".pdf":
        return extract_pdf_text_with_ocr_fallback(path)
    if ext == ".docx":
        try:
            from docx import Document

            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as exc:  # noqa: BLE001 - broad-exception-boundary: document extraction degrades gracefully per file type.
            return f"[DOCX extraction unavailable: {exc}]"
    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook

            wb = load_workbook(path, read_only=True, data_only=True)
            lines: list[str] = []
            for ws in wb.worksheets:
                lines.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    lines.append(",".join("" if value is None else str(value) for value in row))
            return "\n".join(lines)
        except Exception as exc:  # noqa: BLE001 - broad-exception-boundary: document extraction degrades gracefully per file type.
            return f"[XLSX extraction unavailable: {exc}]"
    if ext == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            lines = []
            for idx, slide in enumerate(prs.slides, start=1):
                lines.append(f"# Slide {idx}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            return "\n".join(lines)
        except Exception as exc:  # noqa: BLE001 - broad-exception-boundary: document extraction degrades gracefully per file type.
            return f"[PPTX extraction unavailable: {exc}]"
    return "[Unsupported document type]"


def _record_document_output(
    context: dict[str, Any],
    output: dict[str, Any],
    *,
    source_content: Any,
    source_id: str,
    field_lineage: list[dict[str, str]],
    source_kind: str = "document",
    origin: str = "local_document",
) -> dict[str, Any]:
    record_tool_output_provenance(
        context,
        output,
        source_content=source_content,
        source_kind=source_kind,
        source_id=source_id,
        origin=origin,
        trust_level="untrusted",
        taint_flags=["external_content", "document_content"],
        task_scope=str(context.get("task_id") or ""),
        field_lineage=field_lineage,
    )
    return output


@_require_document_ai_tool
def extract_text(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    extracted = text[:_EXTRACT_TEXT_LIMIT]
    output = {
        "path": str(path),
        "text": extracted,
        "truncated": len(text) > _EXTRACT_TEXT_LIMIT,
    }
    return _record_document_output(
        context,
        output,
        source_content=extracted,
        source_id=str(path),
        field_lineage=[
            {"output_pointer": "/text", "source_pointer": "", "operation": "extract"},
        ],
    )


@_require_document_ai_tool
def summarize(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    result = document_service.summarize_text(
        text,
        path_label=path.name,
        max_chars_to_llm=_document_max_chars_to_llm(context),
        provider_resolver=_document_provider(context),
    )
    output = {"path": str(path), **result}
    return _record_document_output(
        context,
        output,
        source_content=text,
        source_id=str(path),
        field_lineage=[
            {"output_pointer": "/summary", "source_pointer": "", "operation": "summarize"},
        ],
    )


@_require_document_ai_tool
def qa(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    question = str(args.get("question") or "").strip()
    text = extract_text_from_path(path)
    result = document_service.answer_question(
        text,
        question,
        path_label=path.name,
        max_chars_to_llm=_document_max_chars_to_llm(context),
        provider_resolver=_document_provider(context),
    )
    output = {"path": str(path), **result}
    return _record_document_output(
        context,
        output,
        source_content=text,
        source_id=str(path),
        field_lineage=[
            {"output_pointer": "/answer", "source_pointer": "", "operation": "rewrite"},
        ],
    )


@_require_document_ai_tool
def convert_to_markdown(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    output = {"markdown": f"# {path.name}\n\n{text}"[:_EXTRACT_TEXT_LIMIT]}
    return _record_document_output(
        context,
        output,
        source_content=text,
        source_id=str(path),
        field_lineage=[
            {"output_pointer": "/markdown", "source_pointer": "", "operation": "rewrite"},
        ],
    )


@_require_document_ai_tool
def analyze_csv(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    document_intelligence_service._ensure_parseable_file_size(path)
    row_limit = _max_csv_analysis_rows()
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.DictReader(handle)
        columns: list[str | None] = []
        row_count = 0
        for row in reader:
            row_count += 1
            if row_count > row_limit:
                raise document_intelligence_service.DocumentTooLargeError(
                    f"CSV exceeds row limit ({row_count} rows; max {row_limit})."
                )
            if row_count == 1:
                columns = list(row.keys())
    return {"path": str(path), "rows": row_count, "columns": columns}


@_require_document_ai_tool
def analyze_xlsx(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    text = extract_text_from_path(path)
    preview = text[:2000]
    output = {"path": str(path), "preview": preview}
    return _record_document_output(
        context,
        output,
        source_content=preview,
        source_id=str(path),
        field_lineage=[
            {"output_pointer": "/preview", "source_pointer": "", "operation": "extract"},
        ],
    )


@_require_document_ai_tool
def generate_report(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    content = str(args.get("content") or "").strip()
    title = str(args.get("title") or "Report").strip() or "Report"
    output = document_service.generate_report(
        content,
        title=title,
        max_chars_to_llm=_document_max_chars_to_llm(context),
        provider_resolver=_document_provider(context),
    )
    return _record_document_output(
        context,
        output,
        source_content=content,
        source_id=str(context.get("step_id") or context.get("task_id") or "document.generate_report:input"),
        source_kind="document_input",
        origin="tool_argument",
        field_lineage=[
            {"output_pointer": "/report", "source_pointer": "", "operation": "rewrite"},
        ],
    )


@_require_document_ai_tool
def parse_advanced(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    result = document_intelligence_service.parse_advanced(path, settings=context.get("settings"))
    return result.as_dict()


@_require_document_ai_tool
def extract_tables(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    output = document_intelligence_service.extract_tables(path, settings=context.get("settings"))
    tables = output.get("tables") if isinstance(output.get("tables"), list) else []
    return _record_document_output(
        context,
        output,
        source_content=tables,
        source_id=str(path),
        source_kind="document_ir",
        origin="document_parser",
        field_lineage=[
            {"output_pointer": "/tables", "source_pointer": "", "operation": "extract"},
        ],
    )


@_require_document_ai_tool
def ask_with_citations(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    output = document_intelligence_service.ask_with_citations(
        path,
        str(args.get("question") or ""),
        settings=context.get("settings"),
        provider_resolver=_document_provider(context),
        top_k=int(args.get("top_k") or document_intelligence_service.DEFAULT_TOP_K),
    )
    source_blocks = output.get("source_blocks") if isinstance(output.get("source_blocks"), list) else []
    return _record_document_output(
        context,
        output,
        source_content=source_blocks,
        source_id=str(path),
        field_lineage=[
            {"output_pointer": "/answer", "source_pointer": "", "operation": "rewrite"},
        ],
    )


@_require_document_ai_tool
def compare(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    left_raw, right_raw = _compare_path_pair(args)
    left_path = resolve_authorized(left_raw, _allowed(context))
    right_path = resolve_authorized(right_raw, _allowed(context))
    return document_intelligence_service.compare_documents(
        left_path,
        right_path,
        settings=context.get("settings"),
    )


@_require_document_ai_tool
def redact_preview(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    custom_patterns = args.get("custom_patterns") if isinstance(args.get("custom_patterns"), dict) else None
    return document_intelligence_service.redact_preview(
        path,
        settings=context.get("settings"),
        custom_patterns=custom_patterns,
        max_chars=int(args.get("max_chars") or document_intelligence_service.DEFAULT_PREVIEW_CHARS),
    )


@_require_document_ai_tool
def generate_cited_report(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    output = document_intelligence_service.generate_cited_report(
        path,
        title=str(args.get("title") or "Cited Report"),
        query=str(args.get("query") or ""),
        settings=context.get("settings"),
        provider_resolver=_document_provider(context),
        max_blocks=int(args.get("max_blocks") or document_intelligence_service.DEFAULT_REPORT_BLOCKS),
    )
    source_blocks = output.get("source_blocks") if isinstance(output.get("source_blocks"), list) else []
    return _record_document_output(
        context,
        output,
        source_content=source_blocks,
        source_id=str(path),
        field_lineage=[
            {"output_pointer": "/report", "source_pointer": "", "operation": "summarize"},
        ],
    )


@_require_document_ai_tool
def edit_docx(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    dry_run = bool(args.get("dry_run", True))
    if not dry_run:
        raise_if_tool_aborted(context)
    return document_intelligence_service.edit_docx(
        path,
        find=str(args.get("find") or ""),
        replace=str(args.get("replace") or ""),
        dry_run=dry_run,
        abort_context=context,
    )


@_require_document_ai_tool
def edit_xlsx(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    dry_run = bool(args.get("dry_run", True))
    if not dry_run:
        raise_if_tool_aborted(context)
    return document_intelligence_service.edit_xlsx(
        path,
        sheet=str(args.get("sheet") or ""),
        cell=str(args.get("cell") or ""),
        value=args.get("value"),
        dry_run=dry_run,
        abort_context=context,
    )


@_require_document_ai_tool
def edit_pptx(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    dry_run = bool(args.get("dry_run", True))
    if not dry_run:
        raise_if_tool_aborted(context)
    return document_intelligence_service.edit_pptx(
        path,
        find=str(args.get("find") or ""),
        replace=str(args.get("replace") or ""),
        dry_run=dry_run,
        abort_context=context,
    )


@_require_document_ai_tool
def apply_redaction(args: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    path = resolve_authorized(args["path"], _allowed(context))
    custom_patterns = args.get("custom_patterns") if isinstance(args.get("custom_patterns"), dict) else None
    dry_run = bool(args.get("dry_run", True))
    if not dry_run:
        raise_if_tool_aborted(context)
    return document_intelligence_service.apply_redaction(
        path,
        settings=context.get("settings"),
        custom_patterns=custom_patterns,
        max_chars=int(args.get("max_chars") or document_intelligence_service.DEFAULT_PREVIEW_CHARS),
        dry_run=dry_run,
        abort_context=context,
    )


def _compare_path_pair(args: dict[str, Any]) -> tuple[Any, Any]:
    paths = args.get("paths") if isinstance(args.get("paths"), list) else []
    left = args.get("left_path") or (paths[0] if len(paths) > 0 else None)
    right = args.get("right_path") or (paths[1] if len(paths) > 1 else None)
    return left, right


def _validate_compare(args: dict[str, Any], context: dict[str, Any]) -> None:  # noqa: ARG001
    left, right = _compare_path_pair(args)
    if not left or not right:
        raise ValueError(
            "document.compare needs two document paths: provide both left_path and right_path, "
            "or a two-element paths list."
        )
    if not isinstance(left, str) or not isinstance(right, str):
        raise ValueError("document.compare paths must be file path strings.")


def _validate_question(args: dict[str, Any], context: dict[str, Any]) -> None:  # noqa: ARG001
    if not str(args.get("question") or "").strip():
        raise ValueError("This tool needs a non-empty 'question' to answer from the document.")


def _validate_edit_pptx(args: dict[str, Any], context: dict[str, Any]) -> None:  # noqa: ARG001
    if not str(args.get("find") or "").strip():
        raise ValueError("document.edit_pptx requires a non-empty 'find' string.")


def _validate_edit_docx(args: dict[str, Any], context: dict[str, Any]) -> None:  # noqa: ARG001
    if not str(args.get("find") or "").strip():
        raise ValueError("document.edit_docx requires a non-empty 'find' string.")


def _validate_edit_xlsx(args: dict[str, Any], context: dict[str, Any]) -> None:  # noqa: ARG001
    if not str(args.get("sheet") or "").strip() or not str(args.get("cell") or "").strip():
        raise ValueError("document.edit_xlsx requires non-empty 'sheet' and 'cell'.")


def register(registry) -> None:
    path_prop = {"type": "string", "description": "Authorized path to the document to read."}

    def _path_schema(extra: dict[str, Any] | None = None, *, required: list[str] | None = None) -> dict[str, Any]:
        props: dict[str, Any] = {"path": path_prop}
        if extra:
            props.update(extra)
        return {"type": "object", "properties": props, "required": required if required is not None else ["path"]}

    schemas: dict[str, dict[str, Any]] = {
        "document.extract_text": _path_schema(),
        "document.summarize": _path_schema(),
        "document.qa": _path_schema(
            {"question": {"type": "string", "description": "Question to answer from the document."}},
            required=["path", "question"],
        ),
        "document.convert_to_markdown": _path_schema(),
        "document.analyze_csv": _path_schema(),
        "document.analyze_xlsx": _path_schema(),
        "document.generate_report": {
            "type": "object",
            "properties": {
                "content": {"type": "string", "description": "Source content to turn into a report."},
                "title": {"type": "string", "description": "Optional report title."},
            },
            "required": ["content"],
        },
        "document.parse_advanced": _path_schema(),
        "document.extract_tables": _path_schema(),
        "document.ask_with_citations": _path_schema(
            {
                "question": {"type": "string", "description": "Question to answer with cited evidence."},
                "top_k": {"type": "integer", "description": "Optional number of evidence blocks to retrieve."},
            },
            required=["path", "question"],
        ),
        "document.compare": {
            "type": "object",
            "properties": {
                "left_path": {"type": "string", "description": "First document path."},
                "right_path": {"type": "string", "description": "Second document path."},
                "paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Alternative two-element list of document paths.",
                },
            },
            "required": [],
        },
        "document.redact_preview": _path_schema(
            {
                "custom_patterns": {"type": "object", "description": "Optional custom redaction patterns."},
                "max_chars": {"type": "integer", "description": "Optional preview character cap."},
            }
        ),
        "document.generate_cited_report": _path_schema(
            {
                "title": {"type": "string", "description": "Optional report title."},
                "query": {"type": "string", "description": "Optional focusing query for the report."},
                "max_blocks": {"type": "integer", "description": "Optional cap on cited blocks."},
            }
        ),
        "document.edit_docx": _path_schema(
            {
                "find": {"type": "string", "description": "Text to find in the document."},
                "replace": {"type": "string", "description": "Replacement text."},
                "dry_run": {"type": "boolean", "description": "Preview matches without writing (default true)."},
            },
            required=["path", "find", "replace"],
        ),
        "document.edit_xlsx": _path_schema(
            {
                "sheet": {"type": "string", "description": "Worksheet name."},
                "cell": {"type": "string", "description": "Cell reference such as B2."},
                "value": {"description": "New cell value."},
                "dry_run": {"type": "boolean", "description": "Preview change without writing (default true)."},
            },
            required=["path", "sheet", "cell", "value"],
        ),
        "document.edit_pptx": _path_schema(
            {
                "find": {"type": "string", "description": "Text to find in slide shapes."},
                "replace": {"type": "string", "description": "Replacement text."},
                "dry_run": {"type": "boolean", "description": "Preview matches without writing (default true)."},
            },
            required=["path", "find", "replace"],
        ),
        "document.apply_redaction": _path_schema(
            {
                "custom_patterns": {"type": "object", "description": "Optional custom redaction patterns."},
                "max_chars": {"type": "integer", "description": "Optional preview character cap."},
                "dry_run": {"type": "boolean", "description": "Preview redaction without writing (default true)."},
            }
        ),
    }

    write_defs = [
        ("document.edit_docx", edit_docx, _validate_edit_docx),
        ("document.edit_xlsx", edit_xlsx, _validate_edit_xlsx),
        ("document.edit_pptx", edit_pptx, _validate_edit_pptx),
        ("document.apply_redaction", apply_redaction, None),
    ]
    read_defs = [
        ("document.extract_text", extract_text),
        ("document.summarize", summarize),
        ("document.qa", qa),
        ("document.convert_to_markdown", convert_to_markdown),
        ("document.analyze_csv", analyze_csv),
        ("document.analyze_xlsx", analyze_xlsx),
        ("document.generate_report", generate_report),
        ("document.parse_advanced", parse_advanced),
        ("document.extract_tables", extract_tables),
        ("document.ask_with_citations", ask_with_citations),
        ("document.compare", compare),
        ("document.redact_preview", redact_preview),
        ("document.generate_cited_report", generate_cited_report),
    ]
    defs = read_defs
    validators = {
        "document.qa": _validate_question,
        "document.ask_with_citations": _validate_question,
        "document.compare": _validate_compare,
        "document.edit_docx": _validate_edit_docx,
        "document.edit_xlsx": _validate_edit_xlsx,
        "document.edit_pptx": _validate_edit_pptx,
    }
    for name, fn in defs:
        registry.register(
            ToolDefinition(
                name=name,
                description=tool_description(name),
                search_hint=tool_search_hint(name),
                input_schema=schemas[name],
                output_schema={},
                risk_level=RiskLevel.R0_READ_ONLY,
                agent_owner="DocumentAgent",
                supports_dry_run=False,
                requires_authorized_path=name != "document.generate_report",
                execute=fn,
                validate_input=validators.get(name),
                read_only=True,
                concurrency_safe=True,
                effects=["read"],
                resource_kinds=["document"],
                fast_path_eligible=True,
                trust_tier="builtin",
            )
        )
    write_output_schema = {
        "type": "object",
        "properties": {
            "ok": {"type": "boolean"},
            "dry_run": {"type": "boolean"},
            "path": {"type": "string"},
            "changed_paths": {"type": "array", "items": {"type": "string"}},
            "diff_preview": {"type": "array"},
            "rollback_info": {"type": "object"},
            "match_count": {"type": "integer"},
            "error_code": {"type": "string"},
        },
    }
    for name, fn, validator in write_defs:
        registry.register(
            ToolDefinition(
                name=name,
                description=tool_description(name),
                search_hint=tool_search_hint(name),
                input_schema=schemas[name],
                output_schema=write_output_schema,
                risk_level=RiskLevel.R2_REVERSIBLE_MODIFY,
                agent_owner="DocumentAgent",
                supports_dry_run=True,
                requires_authorized_path=True,
                execute=fn,
                validate_input=validator,
                read_only=False,
                concurrency_safe=False,
                effects=["write"],
                resource_kinds=["document"],
                fast_path_eligible=False,
                trust_tier="builtin",
            )
        )
